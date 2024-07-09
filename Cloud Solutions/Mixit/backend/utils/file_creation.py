import fitz as pyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from spire.presentation import *

from utils.docx_file_helper import docx_to_pdf


def create_pptx(title_text: str, subtitle_text: str):
    prs = PptxPresentation("assets/powerpoint_template.pptx")

    first_slide = prs.slides[0]
    title = first_slide.shapes.title
    subtitle = first_slide.placeholders[1]

    title.text = title_text
    subtitle.text = subtitle_text

    file_name = f"{title_text}.pptx"
    prs.save(file_name)
    return file_name


def create_docx(title_text: str, subtitle_text: str):
    doc = DocxDocument("assets/word_template.docx")
    doc.paragraphs[6].add_run(title_text)
    doc.paragraphs[7].add_run(subtitle_text)

    file_name = f"{title_text}.docx"
    doc.save(file_name)
    return file_name


def pptx_to_pdf(pptx_file: str, only_first_slide: bool = True):
    presentation = Presentation()
    presentation.LoadFromFile(pptx_file)
    slides = presentation
    if only_first_slide:
        slide = presentation.Slides[0]
    slides.SaveToFile("SlideToPDF.pdf", FileFormat.PDF)
    presentation.Dispose()
    return "SlideToPDF.pdf"


def pdf_to_images(pdf_file: str):
    doc = pyMuPDF.open(pdf_file)
    image = doc[0].get_pixmap()
    image.save("SlideToPDF.png")
    doc.close()
    return "SlideToPDF.png"


def get_preview_image(file_name: str):
    if file_name.endswith('.pptx'):
        return pdf_to_images(pptx_to_pdf(file_name))
    elif file_name.endswith('.docx'):
        return pdf_to_images(docx_to_pdf(file_name))
    return None
